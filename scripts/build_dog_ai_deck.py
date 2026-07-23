#!/usr/bin/env python3
"""강아지 AI 에이전트 산업 조사 발표자료 생성기 (python-pptx -> Keynote 호환 pptx)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 팔레트 ----
INK = RGBColor(0x16, 0x26, 0x3A)
INK2 = RGBColor(0x2B, 0x3C, 0x55)
ACCENT = RGBColor(0xFF, 0x6B, 0x4A)      # coral
ACCENT2 = RGBColor(0x1F, 0xA9, 0x8F)     # teal
WARN = RGBColor(0xC2, 0x5A, 0x3F)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LINE = RGBColor(0xD8, 0xDC, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF6, 0xF4, 0xEF)
LIGHTBG = RGBColor(0xF7, 0xF6, 0xF2)

FONT = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def style(run, size, bold=False, color=INK, name=FONT):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.name = name
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.shadow.inherit = False
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return r


def rect(s, x, y, w, h, fill=None, line=None, line_w=None, rounded=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = s.shapes.add_shape(shp, x, y, w, h)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = line_w or Pt(1)
    return sh


def tbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def indent(p, marl, ind=0):
    pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    pPr.set("marL", str(int(marl)))
    pPr.set("indent", str(int(ind)))


def para(tf, runs, first=False, gap=12, line=1.12, align=None, marl=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(gap)
    p.space_before = Pt(0)
    p.line_spacing = line
    if align is not None:
        p.alignment = align
    if marl is not None:
        indent(p, marl)
    for text, size, bold, color in runs:
        r = p.add_run()
        r.text = text
        style(r, size, bold, color)
    return p


def footer(s, idx):
    tf = tbox(s, Inches(0.9), Inches(7.02), Inches(8), Inches(0.35))
    para(tf, [("강아지 AI 에이전트 산업   ·   2026.05", 9, False, MUTED)], first=True, gap=0)
    tf2 = tbox(s, Inches(11.6), Inches(7.0), Inches(0.83), Inches(0.35))
    para(tf2, [(f"{idx:02d}", 11, True, MUTED)], first=True, gap=0, align=PP_ALIGN.RIGHT)


def header(s, kicker, title, idx):
    rect(s, Inches(0.9), Inches(0.66), Inches(0.22), Inches(0.22), fill=ACCENT)
    tf = tbox(s, Inches(1.24), Inches(0.6), Inches(10), Inches(0.4))
    para(tf, [(kicker, 12.5, True, ACCENT)], first=True, gap=0)
    tf2 = tbox(s, Inches(0.9), Inches(1.0), Inches(11.6), Inches(1.0))
    para(tf2, [(title, 29, True, INK)], first=True, gap=0)
    rect(s, Inches(0.9), Inches(1.94), Inches(11.53), Pt(2.4), fill=INK)
    footer(s, idx)


# =================================================================
# S1 — 표지
# =================================================================
s = slide()
bg(s, INK)
rect(s, Inches(0.9), Inches(2.7), Inches(1.1), Inches(0.16), fill=ACCENT)
tf = tbox(s, Inches(0.9), Inches(1.5), Inches(11), Inches(0.5))
para(tf, [("AI 창업 전략 리포트", 14, True, ACCENT2)], first=True, gap=0)
tf = tbox(s, Inches(0.86), Inches(3.0), Inches(11.6), Inches(2.4))
para(tf, [("강아지 AI 에이전트 산업 지형도와", 40, True, WHITE)], first=True, gap=4, line=1.08)
para(tf, [("솔로 창업 진입 전략", 40, True, WHITE)], gap=0, line=1.08)
tf = tbox(s, Inches(0.9), Inches(5.45), Inches(11), Inches(0.5))
para(tf, [("한국 · 글로벌 동시 분석   /   소프트웨어-only MVP 관점", 16, False, RGBColor(0xB9, 0xC2, 0xD0))], first=True, gap=0)
tf = tbox(s, Inches(0.9), Inches(6.7), Inches(11), Inches(0.4))
para(tf, [("2026.05.28", 12, True, MUTED)], first=True, gap=0)

# =================================================================
# S2 — 한 줄 결론
# =================================================================
s = slide()
header(s, "EXECUTIVE SUMMARY", "한 줄 결론", 2)
tf = tbox(s, Inches(1.0), Inches(2.45), Inches(11.4), Inches(4.2))
para(tf, [("글로벌은 ", 19, False, INK), ("수의 워크플로우 AI·펫보험 AI", 19, True, INK),
          ("에 자금 집중 — 이미 포화·M&A 단계", 19, False, INK)], first=True, gap=18)
para(tf, [("소비자용 ", 19, False, INK), ("짖음 번역·펫 초상화", 19, True, INK),
          ("는 진입 쉽지만 LLM 래퍼 레드오션", 19, False, INK)], gap=18)
para(tf, [("1인 기술창업자의 빈칸 = 대화로 끝까지 따라가며 ", 19, False, INK),
          ("기억하는 에이전트", 19, True, ACCENT)], gap=8)
para(tf, [("사료·영양 어드바이저  ·  반려인 행정 비서  ·  (한국) 한국어 케어 코파일럿", 16, False, MUTED)],
     gap=0, marl=Inches(0.42))

# =================================================================
# S3 — 시장 (글로벌)
# =================================================================
s = slide()
header(s, "MARKET · GLOBAL", "시장 현황 — 글로벌", 3)
tf = tbox(s, Inches(1.0), Inches(2.4), Inches(11.4), Inches(4.2))
items = [
    ("2025년 펫·수의 스타트업 VC 투자 ", "660M 달러+", " 유입 (전년 보합)"),
    ("자금 집중축: ", "수의 AI 스크라이브 · 텔레헬스/트리아지 · 펫보험 AI", ""),
    ("대기업 참전: ", "Mars 펫 헬스 10억 달러", " · Lemonade 청구 55% AI 완전 자동화"),
    ("지역: ", "미국 압도", " · 유럽 Lupa·Invoxia · 중국 PETKIT(하드웨어)"),
    ("성숙 신호: ", "ScribbleVet → Instinct 인수", " (2025)"),
]
first = True
for a, b, c in items:
    para(tf, [("—  ", 18, True, ACCENT), (a, 18, False, INK), (b, 18, True, INK), (c, 18, False, INK)],
         first=first, gap=16)
    first = False

# =================================================================
# S4 — 시장 (한국)
# =================================================================
s = slide()
header(s, "MARKET · KOREA", "시장 현황 — 한국", 4)
tf = tbox(s, Inches(1.0), Inches(2.4), Inches(11.4), Inches(4.2))
items = [
    ("반려가구 ", "591만 (26.7%)", " · 반려인 1,546만 명 (KB 2025)"),
    ("반려견 546만 마리 — ", "사상 첫 감소 (-1.7%)", " / 반려묘 +9.2%"),
    ("월평균 지출 ", "19만 원 (+4만 원)", " : 개체수가 아니라 마리당 깊은 지출이 성장 동력"),
    ("펫보험 침투율 ", "여전히 1.8%", " (저침투·고성장)"),
    ("시장규모 전망 편차: ", "농식품부 21조(2032) ~ KREI 7.6조", " (보수적)"),
]
first = True
for a, b, c in items:
    para(tf, [("—  ", 18, True, ACCENT), (a, 18, False, INK), (b, 18, True, INK), (c, 18, False, INK)],
         first=first, gap=16)
    first = False

# =================================================================
# S5 — 산업 지형도 (표)
# =================================================================
s = slide()
header(s, "LANDSCAPE", "산업 지형도 — 7대 카테고리 온도계", 5)
rows = [
    ("카테고리", "온도", "대표 플레이어"),
    ("수의 클리닉 AI 스크라이브", "과열 · M&A", "ScribbleVet · Digitail · HappyDoc"),
    ("펫보험 + AI 청구", "HOT (대형사)", "Lemonade · Trupanion · 메리츠 펫퍼민트"),
    ("헬스·수의 진단 AI", "HOT", "SignalPET · Lupa · TTcare · 라이펫 · 핏펫"),
    ("웨어러블·행동", "성숙 · 포화", "Fi · Whistle · Invoxia · 펫펄스"),
    ("훈련·짖음 번역", "과열 · 휘발", "Traini · GoodPup · LunaDogAI"),
    ("콘텐츠·추모", "저자본 난립", "DreamPets · Pet2AI · 포포즈"),
    ("스마트 피더·IoT", "자본 필요", "PETKIT · AI Tails"),
]
gfx = s.shapes.add_table(len(rows), 3, Inches(0.9), Inches(2.3), Inches(11.53), Inches(4.3))
tbl = gfx.table
tbl.first_row = False
tbl.horz_banding = False
tbl.columns[0].width = Inches(3.9)
tbl.columns[1].width = Inches(2.4)
tbl.columns[2].width = Inches(5.23)
HOTC = {"과열 · M&A": WARN, "HOT (대형사)": ACCENT, "HOT": ACCENT, "과열 · 휘발": WARN}
for i, row in enumerate(rows):
    tbl.rows[i].height = Inches(0.53)
    for j, val in enumerate(row):
        c = tbl.cell(i, j)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.16)
        c.margin_right = Inches(0.1)
        c.margin_top = Inches(0.02)
        c.margin_bottom = Inches(0.02)
        c.fill.solid()
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = val
        if i == 0:
            c.fill.fore_color.rgb = INK
            style(r, 14, True, WHITE)
        else:
            c.fill.fore_color.rgb = WHITE if i % 2 else LIGHTBG
            if j == 1:
                style(r, 13, True, HOTC.get(val, MUTED))
            elif j == 0:
                style(r, 13.5, True, INK)
            else:
                style(r, 12.5, False, INK2)

# =================================================================
# S6 — 1인 기술창업자 필터 (2단)
# =================================================================
s = slide()
header(s, "FILTER", "1인 기술창업자 필터", 6)
# 왼쪽: 가능
lx = Inches(0.9)
rect(s, lx, Inches(2.35), Inches(5.6), Inches(4.35), fill=LIGHTBG, rounded=True)
rect(s, lx, Inches(2.35), Inches(0.14), Inches(4.35), fill=ACCENT2)
tf = tbox(s, Inches(1.2), Inches(2.62), Inches(5.0), Inches(0.7))
para(tf, [("가능", 19, True, ACCENT2), ("  · 소프트웨어-only LLM 에이전트", 13, True, MUTED)], first=True, gap=0)
tf = tbox(s, Inches(1.2), Inches(3.35), Inches(5.0), Inches(3.1))
for i, t in enumerate(["사료·영양 어드바이저", "반려인 행정 비서", "비진단 트리아지 챗봇",
                       "펫 콘텐츠·추모", "클리닉 전화 AI"]):
    para(tf, [("—  ", 16, True, ACCENT2), (t, 16, False, INK)], first=(i == 0), gap=14)
# 오른쪽: 부적합
rxv = Inches(6.83)
rect(s, rxv, Inches(2.35), Inches(5.6), Inches(4.35), fill=CARD, rounded=True)
rect(s, rxv, Inches(2.35), Inches(0.14), Inches(4.35), fill=WARN)
tf = tbox(s, Inches(7.13), Inches(2.62), Inches(5.0), Inches(0.7))
para(tf, [("부적합", 19, True, WARN), ("  · 하드웨어·규제·자본·선점", 13, True, MUTED)], first=True, gap=0)
tf = tbox(s, Inches(7.13), Inches(3.35), Inches(5.1), Inches(3.1))
pairs = [("웨어러블·스마트피더", "생산·인증"), ("진단영상 AI", "임상검증·의료기기"),
         ("펫보험", "라이선스·자본"), ("분실견 매칭", "Petco 무료 선점"),
         ("입양 매칭", "Amazon 진입")]
for i, (t, why) in enumerate(pairs):
    para(tf, [("—  ", 16, True, WARN), (t, 16, False, INK), (f"  ({why})", 13, False, MUTED)],
         first=(i == 0), gap=14)

# =================================================================
# S7 — 공백 분석
# =================================================================
s = slide()
header(s, "WHITE SPACE · KOREA", "공백 분석 — 한국 LLM 에이전트 빈칸", 7)
tf = tbox(s, Inches(1.0), Inches(2.4), Inches(11.4), Inches(0.7))
para(tf, [("한국은 ", 16, False, MUTED), ("사진 1장 단발 AI 진단", 16, True, INK2),
          ("에 집중 (TTcare·라이펫·핏펫 3파전) → 대화형 종단 에이전트가 비어 있음", 16, False, MUTED)],
     first=True, gap=0)
tf = tbox(s, Inches(1.0), Inches(3.25), Inches(11.4), Inches(3.4))
gaps = [
    ("빈칸 1", "한국어 멀티턴 케어 코파일럿", "문진 대화 → 트리아지 + 영수증·검사지 PDF를 쉬운 한국어로 해설"),
    ("빈칸 2", "노령견·만성질환 케어 매니저", "투약·체중·식이 추세 알림 + 수의사 공유 리포트 자동화"),
    ("빈칸 3", "EMR 위 얹는 AI 애드온", "진료음성 받아쓰기 · 보호자 안내문 생성 (벳칭·아이엠디티 연동)"),
    ("빈칸 4", "펫로스 메모리 에이전트", "사진·영상·대화 로그로 디지털 메모리얼 — 정서 케어"),
]
first = True
for tag, title, desc in gaps:
    para(tf, [(f"{tag}   ", 15, True, ACCENT2), (title, 17, True, INK)], first=first, gap=2)
    para(tf, [(desc, 14, False, MUTED)], gap=14, marl=Inches(0.95))
    first = False

# =================================================================
# S8 — 9개 컨셉 평가표
# =================================================================
s = slide()
header(s, "SCORING", "9개 컨셉 평가표 (솔로 기준)", 8)
rows = [
    ("컨셉", "구현성", "포화", "수익화", "규제위험"),
    ("사료·영양 어드바이저", "높음", "중", "높음", "중"),
    ("반려인 행정 비서 (OCR)", "높음", "중", "중", "낮음"),
    ("펫 콘텐츠·추모", "높음", "높음", "높음", "낮음"),
    ("훈련·행동 코치", "중", "높음", "높음", "중"),
    ("수의 트리아지", "중", "높음", "중", "높음"),
    ("펫 컴패니언(대화)", "높음", "높음", "중", "낮음"),
    ("병원 스크라이브 B2B", "낮음", "높음", "높음", "중"),
    ("분실견 매칭", "낮음", "중", "낮음", "낮음"),
    ("입양 매칭", "중", "중", "낮음", "낮음"),
]
gfx = s.shapes.add_table(len(rows), 5, Inches(0.9), Inches(2.25), Inches(11.53), Inches(4.45))
tbl = gfx.table
tbl.first_row = False
tbl.horz_banding = False
tbl.columns[0].width = Inches(4.33)
for k in range(1, 5):
    tbl.columns[k].width = Inches(1.8)
# 좋음/나쁨 색: 구현성·수익화는 높음=좋음(teal), 포화·규제는 높음=나쁨(warn)
GOOD = ACCENT2
BAD = WARN
MIDC = RGBColor(0xB0, 0x82, 0x2A)
for i, row in enumerate(rows):
    tbl.rows[i].height = Inches(0.44)
    for j, val in enumerate(row):
        c = tbl.cell(i, j)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.16)
        c.margin_right = Inches(0.05)
        c.margin_top = Inches(0.01)
        c.margin_bottom = Inches(0.01)
        c.fill.solid()
        p = c.text_frame.paragraphs[0]
        if j >= 1:
            p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = val
        if i == 0:
            c.fill.fore_color.rgb = INK
            style(r, 13, True, WHITE)
        else:
            c.fill.fore_color.rgb = WHITE if i % 2 else LIGHTBG
            if j == 0:
                style(r, 13, True, INK)
            else:
                positive_high = j in (1, 3)  # 구현성, 수익화
                if val == "높음":
                    col = GOOD if positive_high else BAD
                elif val == "낮음":
                    col = BAD if positive_high else GOOD
                else:
                    col = MIDC
                style(r, 13, True, col)

# =================================================================
# S9 — TOP 3 추천 (3단 카드)
# =================================================================
s = slide()
header(s, "RECOMMENDATION", "TOP 3 추천", 9)
cards = [
    ("01", "사료·영양 어드바이저", ACCENT,
     [("WEDGE", "일회성 플랜이 아닌 체중·재주문·간식칼로리 지속 코칭"),
      ("수익", "구독 + 사료·보충제 제휴 커미션"),
      ("함정", "처방 톤 금지 · AAFCO·수의 가이드 인용 강제")]),
    ("02", "반려인 행정 비서", ACCENT2,
     [("WEDGE", "영수증·접종증명 사진 → 일정·약·접종 캘린더 자동화"),
      ("빈칸", "종이기록 OCR 자동 파싱 내세운 앱 없음"),
      ("함정", "능동 행동 없으면 무료 캘린더에 패배")]),
    ("03", "펫 콘텐츠·추모", INK2,
     [("WEDGE", "감정 결제 즉시 전환 + 바이럴 캐시플로우"),
      ("포지션", "한국 시장 특화 재포지셔닝"),
      ("함정", "모델 래퍼 → 정서 브랜드·커뮤니티가 해자")]),
]
cw = Inches(3.64)
gap = Inches(0.305)
x0 = Inches(0.9)
for k, (rank, name, col, lines) in enumerate(cards):
    x = x0 + (cw + gap) * k
    rect(s, x, Inches(2.35), cw, Inches(4.35), fill=LIGHTBG, rounded=True)
    rect(s, x, Inches(2.35), cw, Inches(0.7), fill=col, rounded=True)
    rect(s, x, Inches(2.7), cw, Inches(0.35), fill=col)  # 모서리 메우기
    tf = tbox(s, x + Inches(0.25), Inches(2.42), Inches(1), Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(rank, 26, True, WHITE)], first=True, gap=0)
    tf = tbox(s, x + Inches(1.0), Inches(2.42), cw - Inches(1.2), Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(name, 16.5, True, WHITE)], first=True, gap=0, align=PP_ALIGN.RIGHT)
    tf = tbox(s, x + Inches(0.28), Inches(3.3), cw - Inches(0.56), Inches(3.2))
    first = True
    for tag, desc in lines:
        para(tf, [(tag, 12, True, col)], first=first, gap=2)
        para(tf, [(desc, 13.5, False, INK)], gap=15, line=1.1)
        first = False

# =================================================================
# S10 — AI 인프라 & 해자
# =================================================================
s = slide()
header(s, "INFRA & MOAT", "AI 인프라 & 해자", 10)
tf = tbox(s, Inches(1.0), Inches(2.4), Inches(11.4), Inches(4.2))
items = [
    ("두뇌: ", "Gemini Flash-Lite / GPT-4o-mini", "로 대화·플랜 충분, 깊은 추론만 Claude Sonnet"),
    ("비전: ", "라벨·증명서·사진 인식", " = 사용자당 월 수십 원대"),
    ("콘텐츠: ", "이미지 1장 ~1,120토큰", ", 객단가 $6~10이면 마진 충분"),
    ("해자는 모델이 아님 → ", "누적 데이터·메모리 / 제휴 네트워크 / 정서 브랜드", ""),
    ("설계 원칙: ", "MVP부터 기억하는 에이전트", " — 전환비용 축적"),
]
first = True
for a, b, c in items:
    para(tf, [("—  ", 18, True, ACCENT), (a, 18, False, INK), (b, 18, True, INK), (c, 18, False, INK)],
         first=first, gap=16)
    first = False

# =================================================================
# S11 — 규제 리스크
# =================================================================
s = slide()
header(s, "REGULATION · KOREA", "규제 리스크 — 한국 핵심 병목", 11)
tf = tbox(s, Inches(1.0), Inches(2.4), Inches(11.4), Inches(4.2))
items = [
    ("수의사법: 대면진료 원칙", " → 비대면은 상담·조언 수준으로 제한"),
    ("에이아이포펫 규제샌드박스 1호 실증특례", " (재진 안과) — 현재 국내 유일"),
    ("전략: 초기엔 진단·처방 표방 금지", " → 정보정리·트리아지·기록·커뮤니케이션 레이어"),
    ("정부 '반려동물 연관산업 육성법' 추진", " → 펫테크 지원사업 활용"),
]
first = True
for b, c in items:
    para(tf, [("—  ", 18, True, ACCENT), (b, 18, True, INK), (c, 18, False, INK)], first=first, gap=18)
    first = False

# =================================================================
# S12 — 90일 로드맵
# =================================================================
s = slide()
header(s, "ROADMAP", "90일 실행 로드맵 (솔로)", 12)
phases = [
    ("0–30일", "영양 MVP — 라벨 사진 → 적합도 판정", "보호자 20명 인터뷰로 문제 검증"),
    ("30–60일", "기억하는 프로필 (견종·체중·알레르기)", "재주문·체중 추적 루프 구축"),
    ("60–90일", "제휴 사료 1~2곳 커미션 연동", "유료 구독 베타 + 리텐션 측정"),
]
y = Inches(2.5)
for i, (ph, title, desc) in enumerate(phases):
    yy = y + Inches(1.18) * i
    rect(s, Inches(0.9), yy, Inches(1.7), Inches(0.92), fill=INK, rounded=True)
    tf = tbox(s, Inches(0.9), yy, Inches(1.7), Inches(0.92), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(ph, 17, True, WHITE)], first=True, gap=0, align=PP_ALIGN.CENTER)
    tf = tbox(s, Inches(2.85), yy + Inches(0.06), Inches(9.5), Inches(0.85), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(title, 17, True, INK)], first=True, gap=3)
    para(tf, [(desc, 14, False, MUTED)], gap=0)
rect(s, Inches(0.9), Inches(6.35), Inches(11.53), Pt(1.2), fill=LINE)
tf = tbox(s, Inches(0.9), Inches(6.5), Inches(11.4), Inches(0.5))
para(tf, [("검증 지표  ", 14, True, ACCENT2), ("7일 리텐션  ·  주간 활성 대화 수  ·  결제 전환율", 14, False, INK)],
     first=True, gap=0)

# =================================================================
# S13 — 결론 & 다음 액션
# =================================================================
s = slide()
bg(s, INK)
rect(s, Inches(0.9), Inches(0.9), Inches(1.1), Inches(0.16), fill=ACCENT)
tf = tbox(s, Inches(0.9), Inches(1.25), Inches(11), Inches(0.9))
para(tf, [("결론 & 다음 액션", 32, True, WHITE)], first=True, gap=0)
tf = tbox(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(4))
concl = [
    ("큰 자금이 몰린 곳", "(스크라이브·보험·진단영상)은 솔로의 싸움터가 아님"),
    ("빈칸 = ", "대화로 끝까지 따라가며 기억하는 한국어 에이전트"),
    ("1순위 사료·영양 → 행정 비서로 확장", " (동일 반려동물 프로필 자산 재사용)"),
    ("규제는 진단 표방 금지로 우회", ", 정서·데이터·제휴로 해자 구축"),
]
first = True
for a, b in concl:
    para(tf, [("—  ", 19, True, ACCENT2), (a, 19, True, WHITE), (b, 19, False, RGBColor(0xC4, 0xCC, 0xD8))],
         first=first, gap=20)
    first = False
footer(s, 13)

import os
out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   "강아지-AI-에이전트-창업-발표.pptx")
prs.save(out)
print(out)
